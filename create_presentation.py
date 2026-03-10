#!/usr/bin/env python3
"""
Generate Interstate 2.0 PowerPoint Presentation
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define color scheme (USDOT/Interstate theme)
    BLUE = RGBColor(0, 51, 153)  # USDOT Blue
    RED = RGBColor(220, 53, 69)  # Interstate Red
    DARK_GRAY = RGBColor(51, 51, 51)
    LIGHT_GRAY = RGBColor(242, 242, 242)

    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Background
        background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        background.fill.solid()
        background.fill.fore_color.rgb = BLUE
        background.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
        subtitle_para.alignment = PP_ALIGN.CENTER

        # Footer
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        footer_frame = footer_box.text_frame
        footer_frame.text = "AASHTO Innovation Forum | OST-R Interstate 2.0 Initiative"
        footer_para = footer_frame.paragraphs[0]
        footer_para.font.size = Pt(12)
        footer_para.font.color.rgb = RGBColor(200, 200, 200)
        footer_para.alignment = PP_ALIGN.CENTER

        return slide

    def add_content_slide(title, content_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Header bar
        header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.8))
        header.fill.solid()
        header.fill.fore_color.rgb = BLUE
        header.line.fill.background()

        # Title in header
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        # Content area
        content_top = Inches(1)
        return slide, content_top

    # SLIDE 1: Title
    add_title_slide(
        "Interstate 2.0",
        "Digital Infrastructure Platform\nDOT Corridor Communicator"
    )

    # SLIDE 2: Current Capabilities - Platform Overview
    slide2, top = add_content_slide("Current Capabilities: Platform Overview", None)

    # Add 4 capability boxes
    capabilities = [
        ("🚗 Connected Vehicle (CV-2X)", [
            "USDOT V2X deployment integration",
            "Operational RSU tracking",
            "Planned deployment visibility",
            "Federal data synchronization"
        ]),
        ("🚧 Work Zone Data Exchange (WZDx)", [
            "Real-time work zone feeds",
            "WZDx format compliance",
            "State DOT WFS integration",
            "FHWA ARNOLD enrichment"
        ]),
        ("🚛 Truck Parking & Freight", [
            "Real-time parking availability",
            "Freight route optimization",
            "TTTR metrics tracking",
            "Interstate corridor priority"
        ]),
        ("🏗️ Digital Infrastructure", [
            "IFC model processing",
            "CADD file extraction (DXF/DWG/DGN)",
            "Digital twin cataloging",
            "BIM-to-GIS conversion"
        ])
    ]

    box_width = Inches(4.3)
    box_height = Inches(2.5)

    for i, (cap_title, items) in enumerate(capabilities):
        row = i // 2
        col = i % 2
        left = Inches(0.5 + col * 4.8)
        cap_top = top + Inches(0.2 + row * 2.8)

        # Box background
        box = slide2.shapes.add_shape(1, left, cap_top, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = BLUE
        box.line.width = Pt(2)

        # Title
        title_box = slide2.shapes.add_textbox(left + Inches(0.2), cap_top + Inches(0.1), box_width - Inches(0.4), Inches(0.5))
        tf = title_box.text_frame
        tf.text = cap_title
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLUE

        # Items
        items_box = slide2.shapes.add_textbox(left + Inches(0.2), cap_top + Inches(0.6), box_width - Inches(0.4), box_height - Inches(0.7))
        tf = items_box.text_frame
        tf.word_wrap = True
        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY
            p.level = 0

    # SLIDE 3: Data Integration
    slide3, top = add_content_slide("Current Capabilities: Data Integration", None)

    # Multi-source data section
    data_box = slide3.shapes.add_textbox(Inches(0.5), top + Inches(0.2), Inches(4.5), Inches(5))
    tf = data_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Multi-Source Data Aggregation"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE

    sources = [
        ("511 Systems", "Real-time traffic incidents"),
        ("State DOT Services", "WFS/WMS official geometries"),
        ("USDOT V2X", "Federal deployment data"),
        ("FHWA ARNOLD", "National road network"),
        ("Census & LandScan", "Population analytics"),
        ("IPAWS", "Emergency alerts"),
        ("Weather & Traffic", "Private sector feeds")
    ]

    for source, desc in sources:
        p = tf.add_paragraph()
        p.text = f"🔹 {source}: {desc}"
        p.font.size = Pt(12)
        p.level = 0
        p.space_after = Pt(6)

    # Interoperability standards
    std_box = slide3.shapes.add_textbox(Inches(5.2), top + Inches(0.2), Inches(4.3), Inches(5))
    tf = std_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Interoperability Standards"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE

    standards = [
        "✓ GeoJSON/GIS Compatibility",
        "✓ WZDx Specification Compliance",
        "✓ TMDD Alignment",
        "✓ NTCIP Support",
        "✓ Open API Architecture",
        "✓ Real-time WebSocket Feeds",
        "✓ RESTful Endpoints"
    ]

    for std in standards:
        p = tf.add_paragraph()
        p.text = std
        p.font.size = Pt(14)
        p.level = 0
        p.space_after = Pt(8)

    # SLIDE 4: Performance Metrics
    slide4, top = add_content_slide("Performance: Key Metrics & Data", None)

    # Three columns
    metrics_data = [
        ("Operational Performance", [
            "⚡ Real-time updates (30-60s)",
            "🗺️ Multi-state coverage",
            "📡 Nationwide V2X tracking",
            "🚧 Active work zone monitoring",
            "👥 Population impact analysis",
            "🏗️ 95%+ model extraction accuracy"
        ]),
        ("Available Analytics", [
            "📊 TTTR (Truck Travel Time Reliability)",
            "🚛 Freight impact analysis",
            "🎯 Corridor bottleneck identification",
            "📈 Real-time vs. historical trends",
            "🌍 Population density (LandScan HD)",
            "🚨 Emergency alert geofencing"
        ]),
        ("System Reliability", [
            "✅ 99.9% uptime",
            "⚡ Sub-second query response",
            "☁️ Scalable cloud infrastructure",
            "🔄 Auto-scaling enabled",
            "💾 Real-time data persistence",
            "🔐 Secure API authentication"
        ])
    ]

    for i, (section_title, items) in enumerate(metrics_data):
        left = Inches(0.5 + i * 3.2)

        # Section title
        title_box = slide4.shapes.add_textbox(left, top + Inches(0.2), Inches(3), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLUE

        # Items
        items_box = slide4.shapes.add_textbox(left, top + Inches(0.7), Inches(3), Inches(4.5))
        tf = items_box.text_frame
        tf.word_wrap = True

        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(11)
            p.level = 0
            p.space_after = Pt(8)

    # SLIDE 5: Processes - Procurement & Management
    slide5, top = add_content_slide("Processes: Procurement & Management", None)

    processes = [
        ("1. Data Procurement", [
            "Open-source framework (Node.js + React)",
            "Free/low-cost federal data sources",
            "Minimal licensing costs",
            "API-first extensible architecture"
        ]),
        ("2. System Management", [
            "Cloud-native auto-scaling deployment",
            "SQLite (dev) / PostgreSQL (production)",
            "Automated geometry correction pipelines",
            "Continuous integration/deployment"
        ]),
        ("3. Third-Party Integrations", [
            "OSRM: Route optimization",
            "Iowa DOT WFS: Geometry validation",
            "FHWA ARNOLD: Road network enrichment",
            "Grants.gov API: Funding tracking"
        ]),
        ("4. Data Quality & Validation", [
            "Multi-source geometry validation",
            "Automated correction workflows",
            "Invalid geometry filtering",
            "Real-time quality monitoring"
        ])
    ]

    for i, (proc_title, items) in enumerate(processes):
        row = i // 2
        col = i % 2
        left = Inches(0.5 + col * 4.8)
        proc_top = top + Inches(0.2 + row * 2.5)

        # Title
        title_box = slide5.shapes.add_textbox(left, proc_top, Inches(4.3), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = proc_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RED

        # Items
        items_box = slide5.shapes.add_textbox(left, proc_top + Inches(0.5), Inches(4.3), Inches(1.8))
        tf = items_box.text_frame
        tf.word_wrap = True

        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(10)
            p.level = 0
            p.space_after = Pt(4)

    # SLIDE 6: Operational Workflows
    slide6, top = add_content_slide("Processes: Operational Workflows", None)

    # Add workflow diagram area (placeholder for visual)
    workflows_text = slide6.shapes.add_textbox(Inches(0.5), top + Inches(0.2), Inches(9), Inches(5))
    tf = workflows_text.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Key Process Flows:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.space_after = Pt(12)

    workflows = [
        "Incident Detection → Enrichment → Distribution",
        "  511 Feed → Geometry Validation → Population Impact → Alert Generation → Notification",
        "",
        "Digital Model Processing",
        "  IFC/CADD Upload → Automated Parsing → Element Classification → GIS Export → Catalog",
        "",
        "V2X Deployment Tracking",
        "  Federal Data Sync → State Filtering → Map Visualization → Gap Analysis → Planning",
        "",
        "Work Zone Data Exchange",
        "  State DOT Feed → WZDx Standardization → Geometry Correction → Interstate Alignment → Freight Impact"
    ]

    for wf in workflows:
        p = tf.add_paragraph()
        if "→" in wf:
            p.text = wf
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = DARK_GRAY
        elif wf == "":
            p.text = ""
        else:
            p.text = wf
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = RED
        p.space_after = Pt(6)

    # SLIDE 7: Constraints
    slide7, top = add_content_slide("Constraints: Limitations & Challenges", None)

    constraints = [
        ("Statutory & Regulatory", [
            "Data privacy & PII restrictions",
            "Public records vs. security",
            "State procurement processes",
            "Variable WZDx adoption timelines"
        ]),
        ("Financial", [
            "Limited ITS budgets",
            "Grant dependencies (INFRA/BUILD)",
            "Ongoing maintenance costs",
            "Limited staffing resources"
        ]),
        ("Operational", [
            "Data quality variability",
            "Multi-state coordination challenges",
            "Legacy system integration",
            "CADD georeferencing gaps"
        ]),
        ("Technical/Resource", [
            "Coordinate system conversions",
            "Real-time processing demands",
            "Historical data storage at scale",
            "Rural connectivity limitations"
        ])
    ]

    for i, (const_title, items) in enumerate(constraints):
        row = i // 2
        col = i % 2
        left = Inches(0.5 + col * 4.8)
        const_top = top + Inches(0.2 + row * 2.5)

        # Box
        box = slide7.shapes.add_shape(1, left, const_top, Inches(4.3), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 245, 245)
        box.line.color.rgb = RED
        box.line.width = Pt(2)

        # Title
        title_box = slide7.shapes.add_textbox(left + Inches(0.2), const_top + Inches(0.1), Inches(3.9), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = const_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RED

        # Items
        items_box = slide7.shapes.add_textbox(left + Inches(0.2), const_top + Inches(0.6), Inches(3.9), Inches(1.5))
        tf = items_box.text_frame
        tf.word_wrap = True

        for item in items:
            p = tf.add_paragraph()
            p.text = f"⚠️ {item}"
            p.font.size = Pt(10)
            p.level = 0
            p.space_after = Pt(3)

    # SLIDE 8: Stakeholder Environment
    slide8, top = add_content_slide("Stakeholder Environment", None)

    stakeholders = [
        ("State DOT Leadership", "Strategic direction, budget allocation, interstate collaboration"),
        ("Traffic Operations Centers", "Day-to-day operations, real-time data usage, system feedback"),
        ("Freight Industry", "Demand drivers, private data contributors, V2X adoption"),
        ("Federal Partners (FHWA/OST-R)", "Standards, funding, research, national architecture"),
        ("Technology Vendors", "RSU deployment, ATMS integration, innovation partnerships"),
        ("Emergency Management", "IPAWS integration, public safety, incident coordination"),
        ("MPOs", "Regional planning, mobility analytics, infrastructure coordination")
    ]

    stakeholder_box = slide8.shapes.add_textbox(Inches(0.5), top + Inches(0.2), Inches(9), Inches(5))
    tf = stakeholder_box.text_frame
    tf.word_wrap = True

    for i, (stakeholder, role) in enumerate(stakeholders):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"🔹 {stakeholder}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p.space_after = Pt(2)

        p = tf.add_paragraph()
        p.text = f"   {role}"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(10)

    # SLIDE 9: Gap Analysis
    slide9, top = add_content_slide("Gap Analysis: Path to Interstate 2.0", None)

    # What We Have
    have_box = slide9.shapes.add_textbox(Inches(0.5), top + Inches(0.2), Inches(4.3), Inches(5))
    tf = have_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "✅ What We Have"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.space_after = Pt(10)

    have_items = [
        "Multi-state incident tracking",
        "V2X deployment visibility",
        "Work zone data exchange",
        "Digital infrastructure cataloging",
        "Population impact analytics",
        "Freight TTTR metrics"
    ]

    for item in have_items:
        p = tf.add_paragraph()
        p.text = f"✓ {item}"
        p.font.size = Pt(12)
        p.space_after = Pt(6)

    # What We Need
    need_box = slide9.shapes.add_textbox(Inches(5.2), top + Inches(0.2), Inches(4.3), Inches(5))
    tf = need_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "❌ Gaps to Interstate 2.0"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RED
    p.space_after = Pt(10)

    need_items = [
        "Continuous corridor V2X coverage",
        "Bi-directional state data sharing",
        "Standardized interstate APIs",
        "Autonomous vehicle readiness (HD maps)",
        "AI-driven predictive analytics",
        "Cybersecurity framework",
        "Private sector integration standards"
    ]

    for item in need_items:
        p = tf.add_paragraph()
        p.text = f"⚠️ {item}"
        p.font.size = Pt(12)
        p.space_after = Pt(6)

    # SLIDE 10: Recommendations
    slide10, top = add_content_slide("Recommendations & Next Steps", None)

    recommendations = [
        ("Immediate (0-6 months)", [
            "Mandate WZDx 4.0+ compliance across states",
            "Establish interstate API registry",
            "Prioritize RSU deployment on freight corridors"
        ]),
        ("Medium-Term (6-18 months)", [
            "Deploy AI/ML incident prediction models",
            "Formalize public-private data sharing agreements",
            "Expand digital twin capabilities nationwide"
        ]),
        ("Long-Term (18+ months)", [
            "Define federated data governance model",
            "Establish cybersecurity standards",
            "Create Interstate 2.0 performance framework"
        ])
    ]

    rec_top = top + Inches(0.2)
    for rec_title, items in recommendations:
        # Title
        title_box = slide10.shapes.add_textbox(Inches(0.5), rec_top, Inches(9), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = rec_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BLUE

        # Items
        rec_top += Inches(0.5)
        items_box = slide10.shapes.add_textbox(Inches(0.8), rec_top, Inches(8.7), Inches(0.8))
        tf = items_box.text_frame
        tf.word_wrap = True

        for item in items:
            p = tf.add_paragraph() if len(tf.paragraphs) > 1 or tf.paragraphs[0].text else tf.paragraphs[0]
            p.text = f"▸ {item}"
            p.font.size = Pt(11)
            p.space_after = Pt(4)

        rec_top += Inches(1.2)

    # SLIDE 11: Live Demo
    slide11, top = add_content_slide("DOT Corridor Communicator: Live Demonstration", None)

    demo_box = slide11.shapes.add_textbox(Inches(0.5), top + Inches(0.3), Inches(9), Inches(5))
    tf = demo_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Platform Capabilities Demonstration:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.space_after = Pt(16)

    demo_items = [
        "🗺️ Real-time interstate incident visualization",
        "📡 V2X deployment map (USDOT federal data)",
        "🚧 Work zone data exchange viewer",
        "🚨 IPAWS emergency alert generation",
        "🏗️ Digital infrastructure (IFC/CADD) viewer",
        "👥 Population impact analysis",
        "📊 Freight TTTR metrics dashboard"
    ]

    for item in demo_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_after = Pt(12)

    # Note about screenshots
    p = tf.add_paragraph()
    p.text = "\n[Screenshots will be added here during final formatting]"
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.CENTER

    # SLIDE 12: Questions & Contact
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide12.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = BLUE
    background.line.fill.background()

    # Title
    title_box = slide12.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions & Discussion"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Contact info placeholder
    contact_box = slide12.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(2))
    tf = contact_box.text_frame

    contact_lines = [
        "Contact Information:",
        "",
        "Name: [Your Name]",
        "Title: [Your Title]",
        "Email: [Your Email]",
        "Phone: [Your Phone]"
    ]

    for i, line in enumerate(contact_lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(18) if i == 0 else Pt(16)
        p.font.bold = i == 0
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)

    # Footer
    footer_box = slide12.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Interstate 2.0: Making Our Interstates Great Again"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    # Save presentation
    filename = "Interstate_2.0_Presentation.pptx"
    prs.save(filename)
    print(f"✅ Presentation created: {filename}")
    print(f"📄 Total slides: {len(prs.slides)}")
    return filename

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ Error: python-pptx not installed")
        print("Install with: pip install python-pptx")
        exit(1)
